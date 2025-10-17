import os
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from dotenv import load_dotenv
from langchain.chains import LLMChain
from langchain_core.prompts import PromptTemplate
from langchain_openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import io
from pydantic import BaseModel

# Load environment variables
load_dotenv()
OPENAI_KEY = os.getenv("OPENAI_API_KEY")
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-3.5-turbo-instruct")

if not OPENAI_KEY:
    raise RuntimeError("OPENAI_API_KEY not set in .env file")

app = FastAPI()

# Middleware for CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- Serve the React build ---
app.mount("/static", StaticFiles(directory="build/static"), name="static")

@app.get("/{full_path:path}")
async def serve_react_app(full_path: str):
    file_path = os.path.join("build", full_path)
    if os.path.exists(file_path) and os.path.isfile(file_path):
        return FileResponse(file_path)
    return FileResponse(os.path.join("build", "index.html"))


# --- AI ROUTE: Generate BMC ---
@app.post("/api/generate_bmc")
async def generate_bmc(file: UploadFile = File(...)):
    try:
        if not file.filename.endswith(".txt"):
            raise HTTPException(status_code=400, detail="Please upload a .txt file")

        # Read file content
        content = (await file.read()).decode("utf-8")

        # Prepare AI prompt
        prompt = PromptTemplate(
            input_variables=["project_description"],
            template=(
                "You are an expert business analyst. Based on the following project description, "
                "generate a **Business Model Canvas** as valid JSON only, Each section should be written in full, descriptive sentences (not bullet points) "
                "Use exactly these keys:\n\n"
                "Key Partners, Key Activities, Value Propositions, Customer Relationships, "
                "Each key must contain multiple points (each on a new line or separated by periods). "
                "Each point must start with an uppercase letter and form a full, meaningful sentence.\n\n"
                "Customer Segments, Key Resources, Channels, Cost Structure, Revenue Streams.\n\n"
                "Project Description:\n{project_description}\n\n"
                "Output ONLY valid JSON"
            )
        )

        llm = OpenAI(temperature=0.3, openai_api_key=OPENAI_KEY, model="gpt-3.5-turbo-instruct")
        chain = LLMChain(prompt=prompt, llm=llm)

        bmc_raw = chain.run({"project_description": content})

        # Try parsing as JSON
        try:
            bmc = json.loads(bmc_raw)
        except Exception:
            bmc = {"raw_output": bmc_raw.strip()}
        
        # Ensure all 9 BMC keys exist
        keys = [
            "Key Partners", "Key Activities", "Value Propositions",
            "Customer Relationships", "Customer Segments",
            "Key Resources", "Channels", "Cost Structure", "Revenue Streams"
        ]    
        for key in keys:
            if key not in bmc:
                bmc[key] = "—"
        
        return {"business_model_canvas": bmc}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# --- PowerPoint Export Model ---
class BMCExportRequest(BaseModel):
    bmc: dict


# --- PowerPoint Export Route ---
@app.post("/api/export_ppt")
async def export_ppt(request: BMCExportRequest):
    try:
        bmc = request.bmc
        
        # Create a presentation
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "The Business Model Canvas"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Define box positions and sizes
        box_top = Inches(1.2)
        box_height_top = Inches(3.5)
        box_height_bottom = Inches(2)
        box_width = Inches(2.4)
        
        # Function to add a box
        def add_box(left, top, width, height, title, content, color):
            box = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, width, height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*color)
            box.line.color.rgb = RGBColor(0, 0, 0)
            box.line.width = Pt(2)
            
            text_frame = box.text_frame
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            text_frame.word_wrap = True
            
            # Add title
            p = text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add content - handle both string and list formats
            if isinstance(content, list):
                for item in content:
                    if item and item.strip():
                        p = text_frame.add_paragraph()
                        p.text = "• " + str(item).strip()
                        p.font.size = Pt(10)
                        p.font.color.rgb = RGBColor(0, 0, 0)
                        p.level = 0
            else:
                content_items = str(content).split('\n') if content else ["—"]
                for item in content_items:
                    if item.strip():
                        p = text_frame.add_paragraph()
                        p.text = "• " + item.strip().capitalize()
                        p.font.size = Pt(10)
                        p.font.color.rgb = RGBColor(0, 0, 0)
                        p.level = 0
        
        # Define light pastel colors for each section
        colors = {
            "Key Partners": (230, 230, 250),           # Lavender
            "Key Activities": (255, 240, 245),         # Lavender Blush
            "Key Resources": (240, 255, 240),          # Honeydew
            "Value Propositions": (255, 250, 205),     # Lemon Chiffon
            "Customer Relationships": (255, 228, 225), # Misty Rose
            "Channels": (240, 248, 255),               # Alice Blue
            "Customer Segments": (255, 245, 238),      # Seashell
            "Cost Structure": (245, 245, 220),         # Beige
            "Revenue Streams": (240, 255, 255)         # Azure
        }
        
        # Top row - Key Partners
        add_box(Inches(0.5), box_top, box_width, box_height_top, 
                "Key Partners", bmc.get("Key Partners", "—"), colors["Key Partners"])
        
        # Key Activities (top half)
        add_box(Inches(3), box_top, box_width, Inches(1.7), 
                "Key Activities", bmc.get("Key Activities", "—"), colors["Key Activities"])
        
        # Key Resources (bottom half)
        add_box(Inches(3), box_top + Inches(1.8), box_width, Inches(1.7), 
                "Key Resources", bmc.get("Key Resources", "—"), colors["Key Resources"])
        
        # Value Propositions (center)
        add_box(Inches(5.5), box_top, box_width, box_height_top, 
                "Value Propositions", bmc.get("Value Propositions", "—"), colors["Value Propositions"])
        
        # Customer Relationships (top half)
        add_box(Inches(8), box_top, box_width, Inches(1.7), 
                "Customer Relationships", bmc.get("Customer Relationships", "—"), colors["Customer Relationships"])
        
        # Channels (bottom half)
        add_box(Inches(8), box_top + Inches(1.8), box_width, Inches(1.7), 
                "Channels", bmc.get("Channels", "—"), colors["Channels"])
        
        # Customer Segments
        add_box(Inches(10.5), box_top, box_width, box_height_top, 
                "Customer Segments", bmc.get("Customer Segments", "—"), colors["Customer Segments"])
        
        # Bottom row
        bottom_top = box_top + box_height_top + Inches(0.1)
        
        # Cost Structure (40%)
        add_box(Inches(0.5), bottom_top, Inches(6), box_height_bottom, 
                "Cost Structure", bmc.get("Cost Structure", "—"), colors["Cost Structure"])
        
        # Revenue Streams (60%)
        add_box(Inches(6.6), bottom_top, Inches(6.3), box_height_bottom, 
                "Revenue Streams", bmc.get("Revenue Streams", "—"), colors["Revenue Streams"])
        
        # Save to BytesIO
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        return StreamingResponse(
            ppt_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=business_model_canvas.pptx"}
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


